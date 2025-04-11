#!/usr/bin/env python3
# -*- coding: utf-8 -*-

import os
from pptx import Presentation
import argparse

def extract_notes(pptx_path, output_path):
    """从PPT文件中提取备注并保存为Markdown格式
    
    Args:
        pptx_path (str): PPT文件路径
        output_path (str): 输出的Markdown文件路径
    """
    # 检查文件是否存在
    if not os.path.exists(pptx_path):
        print(f"错误：找不到PPT文件 '{pptx_path}'")
        return False
    
    try:
        # 打开PPT文件
        prs = Presentation(pptx_path)
        
        # 准备Markdown内容
        markdown_content = f"# {os.path.splitext(os.path.basename(pptx_path))[0]} 备注内容\n\n"
        
        # 遍历所有幻灯片
        for slide_number, slide in enumerate(prs.slides, 1):
            # 获取备注
            notes = slide.notes_slide
            if notes and notes.notes_text_frame.text.strip():
                # 添加幻灯片编号和备注内容
                markdown_content += f"## 第 {slide_number} 张幻灯片\n\n{notes.notes_text_frame.text.strip()}\n\n"
        
        # 保存为Markdown文件
        with open(output_path, 'w', encoding='utf-8') as f:
            f.write(markdown_content)
        
        print(f"备注已成功导出到：{output_path}")
        return True
        
    except Exception as e:
        print(f"处理PPT文件时出错：{str(e)}")
        return False

def main():
    # 设置命令行参数
    parser = argparse.ArgumentParser(description='从PPT文件中提取备注并保存为Markdown格式')
    parser.add_argument('pptx_file', help='PPT文件路径')
    parser.add_argument('-o', '--output', help='输出的Markdown文件路径（可选）')
    
    args = parser.parse_args()
    
    # 如果没有指定输出路径，使用默认路径
    if not args.output:
        output_path = os.path.splitext(args.pptx_file)[0] + '_notes.md'
    else:
        output_path = args.output
    
    # 提取备注
    extract_notes(args.pptx_file, output_path)

if __name__ == '__main__':
    main()